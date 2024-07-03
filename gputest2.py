import streamlit as st
from streamlit_chat import message
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.llms import LlamaCpp
from langchain_community.vectorstores import FAISS
from langchain_community.document_loaders import PyPDFLoader
from langchain.schema import Document as LangchainDocument
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os
import tempfile
from docx import Document as DocxDocument
import openpyxl
from pptx import Presentation
import torch
import asyncio
from datetime import datetime
from googletrans import Translator

# Helper function to initialize a new chat session
def initialize_new_session():
    return {
        "history": [],
        "generated": ["Hello! Ask me anything about this document 🤗"],
        "past": ["Hey! 👋"]
    }

def initialize_session_state():
    if 'sessions' not in st.session_state:
        st.session_state['sessions'] = {}
    if 'current_session' not in st.session_state:
        st.session_state['current_session'] = None
    new_session()  # Create the first session on load

def get_current_session():
    return st.session_state['sessions'].get(st.session_state['current_session'], initialize_new_session())

def switch_session(session_name):
    st.session_state['current_session'] = session_name

def new_session():
    session_name = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    st.session_state['sessions'][session_name] = initialize_new_session()
    switch_session(session_name)

def delete_session(session_name):
    if session_name in st.session_state['sessions']:
        del st.session_state['sessions'][session_name]
        if session_name == st.session_state['current_session']:
            if st.session_state['sessions']:
                st.session_state['current_session'] = list(st.session_state['sessions'].keys())[0]
            else:
                new_session()

def conversation_chat(query, chain):
    session = get_current_session()
    result = chain({"question": query, "chat_history": session['history']})
    session['history'].append((query, result["answer"]))
    
    # Template for a concise response
    concise_answer = result["answer"]  # Use the response directly
    
    return concise_answer

def display_chat_history():
    session = get_current_session()
    reply_container = st.container()
    container = st.container()

    with container:
        with st.form(key='my_form', clear_on_submit=True):
            user_input = st.text_input("Question:", placeholder="Ask about your document", key='input')
            submit_button = st.form_submit_button(label='Send')

        if submit_button and user_input:
            with st.spinner('Generating response...'):
                output = conversation_chat(user_input, st.session_state['chain'])

            session['past'].append(user_input)
            session['generated'].append(output)

    if session['generated']:
        with reply_container:
            for i in range(len(session['generated'])):
                message(session["past"][i], is_user=True, key=str(i) + '_user', avatar_style="thumbs")
                message(session["generated"][i], key=str(i), avatar_style="fun-emoji")

def create_conversational_chain(vector_store):
    # Check GPU availability
    device = 'cuda' if torch.cuda.is_available() else 'cpu'

    # Create llm
    llm = LlamaCpp(
        streaming=True,
        model_path="capybarahermes-2.5-mistral-7b.Q5_K_M.gguf",
        temperature=0.75,
        top_p=1, 
        verbose=True,
        n_ctx=4096,
        use_gpu=True  # Explicitly enable GPU
    )
    
    memory = ConversationBufferMemory(memory_key="chat_history", return_messages=True)

    # Configure the chain to produce more concise responses
    chain = ConversationalRetrievalChain.from_llm(llm=llm, chain_type='stuff',
                                                 retriever=vector_store.as_retriever(search_kwargs={"k": 2}),
                                                 memory=memory)
    return chain

def load_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()

def load_docx(file_path):
    doc = DocxDocument(file_path)
    text = []
    for paragraph in doc.paragraphs:
        text.append(paragraph.text)
    return '\n'.join(text)

def load_xlsx(file_path):
    wb = openpyxl.load_workbook(file_path)
    sheet = wb.active
    text = []
    for row in sheet.iter_rows(values_only=True):
        text.append(' '.join([str(cell) if cell else '' for cell in row]))
    return '\n'.join(text)

def load_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return '\n'.join(text)

async def process_files(uploaded_files):
    documents = []
    for file in uploaded_files:
        file_extension = os.path.splitext(file.name)[1]
        with tempfile.NamedTemporaryFile(delete=False) as temp_file:
            temp_file.write(file.read())
            temp_file_path = temp_file.name

        try:
            if file_extension == ".pdf":
                loader = PyPDFLoader(temp_file_path)
                documents.extend(loader.load())
            elif file_extension == ".txt":
                text_content = load_txt(temp_file_path)
                documents.append(LangchainDocument(page_content=text_content))
            elif file_extension == ".docx":
                text_content = load_docx(temp_file_path)
                documents.append(LangchainDocument(page_content=text_content))
            elif file_extension == ".xlsx":
                text_content = load_xlsx(temp_file_path)
                documents.append(LangchainDocument(page_content=text_content))
            elif file_extension == ".pptx":
                text_content = load_pptx(temp_file_path)
                documents.append(LangchainDocument(page_content=text_content))

        except Exception as e:
            print(f"Error processing {file_extension} file: {e}")

        finally:
            os.remove(temp_file_path)

    return documents

def main():
    # Initialize session state
    initialize_session_state()
    st.title("Document Chat 📄")
    # Initialize Streamlit
    st.sidebar.title("Sessions")
    
    # Create New Session button
    if st.sidebar.button("Create New Session"):
        new_session()

    # Display existing sessions in a dropdown
    sessions = list(st.session_state['sessions'].keys())
    selected_session = st.sidebar.selectbox("Select Session", sessions, index=sessions.index(st.session_state['current_session']) if st.session_state['current_session'] else 0)
    if selected_session != st.session_state['current_session']:
        switch_session(selected_session)

    st.sidebar.title("Upload a Document")
    uploaded_files = st.sidebar.file_uploader("Upload files", accept_multiple_files=True)

    if uploaded_files:
        documents = asyncio.run(process_files(uploaded_files))

        if documents:
            text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=20)
            text_chunks = text_splitter.split_documents(documents)

            # Check for CUDA availability and set device accordingly
            device = 'cuda' if torch.cuda.is_available() else 'cpu'

            # Create embeddings
            embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2", 
                                               model_kwargs={'device': device})

            # Create vector store
            vector_store = FAISS.from_documents(text_chunks, embedding=embeddings)

            # Create the chain object and store it in session state
            st.session_state['chain'] = create_conversational_chain(vector_store)

            display_chat_history()
        
        else:
            st.warning("No documents processed. Please upload a valid file.")
    else:
        display_chat_history()

    # Translation section
    st.sidebar.title("Translation")
    language = st.sidebar.selectbox("Select language", ["en", "es", "fr", "de", "zh-cn"])
    translator = Translator()

    session = get_current_session()
    translated_texts = [translator.translate(text, dest=language).text for text in session['generated']]
    
    with st.container():
        for i, text in enumerate(translated_texts):
            st.write(f"**Response {i+1}:** {text}")

if __name__ == "__main__":
    main()
